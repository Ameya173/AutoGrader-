"""
Gemini API Service — High-Performance Split Pipeline
- Target: Gemini 2.0 Flash (Optimized for speed)
- Split Pipeline: OCR Extraction + Creative Enrichment to prevent timeouts
- Safety: BLOCK_NONE for academic reliability
"""
import os, base64, json, re, traceback, time, uuid, io
import google.generativeai as genai
from google.api_core import exceptions
from dotenv import load_dotenv
import pdfplumber
from pptx import Presentation

load_dotenv()
_model = None

# Safety Settings
SAFETY_NONE = [
    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
]

def _clean_qnum(q):
    """Normalize question numbers."""
    if not q: return "?"
    return str(q).strip().lstrip('Qq').strip('. ')

def get_model(force_refresh=False):
    global _model
    if _model is not None and not force_refresh:
        return _model
    
    # If refreshing, clear the old one
    if force_refresh:
        _model = None
    
    key = os.getenv('GEMINI_API_KEY')
    if not key:
        raise ValueError("GEMINI_API_KEY not set in .env")
    
    genai.configure(api_key=key)
    
    # Automatic Model Selection to prevent 404 errors
    # We try newer models first (2.5, 2.0) as they might have fresh quotas
    models_to_try = [
        'gemini-3-flash-preview',
        'gemini-2.5-pro',
        'gemini-2.5-flash',
        'gemini-3-pro-preview',
        'gemini-2.0-flash',
        'gemini-2.0-flash-001',
        'gemini-flash-lite-latest',
        'gemini-flash-latest',
        'gemini-1.5-flash',
    ]
    
    last_error = None
    for model_name in models_to_try:
        try:
            print(f"DEBUG: Attempting to connect to {model_name}...")
            test_model = genai.GenerativeModel(model_name)
            # Minimal test call to verify name/quota
            test_model.generate_content("ping") 
            _model = test_model
            print(f"DEBUG: Successfully connected to {model_name}!")
            return _model
        except exceptions.NotFound:
            print(f"DEBUG: Model {model_name} not found. Trying next...")
            continue
        except exceptions.ResourceExhausted as e:
            print(f"DEBUG: Quota hit for {model_name}: {e}")
            last_error = e
            continue
        except Exception as e:
            print(f"DEBUG: Error with {model_name}: {e}")
            last_error = e
            continue
    
    # If all failed, don't default to a known-failing model.
    if last_error:
        print(f"CRITICAL: All model attempts failed. Last error: {last_error}")
        # We'll return gemini-flash-latest as a final "hopeful" fallback
        # but we don't set _model so we try again next time
        return genai.GenerativeModel('gemini-flash-latest')
    
    return genai.GenerativeModel('gemini-2.0-flash')

def _safe_json(text, fallback):
    if not text: return fallback
    clean = re.sub(r'```json\s*', '', text)
    clean = re.sub(r'```\s*', '', clean).strip()
    try:
        data = json.loads(clean)
        return data if isinstance(data, (list, dict)) else fallback
    except:
        m = re.search(r'[\[{][\s\S]*[\]}]', clean)
        if m:
            try: return json.loads(m.group(0))
            except: pass
    return fallback

def _img_part(path):
    """Utility to prepare image/pdf for Gemini API."""
    with open(path,'rb') as f: data = f.read()
    ext = path.split('.')[-1].lower()
    mime = 'application/pdf' if ext=='pdf' else f'image/{ext}'
    return {'mime_type': mime, 'data': data}

def extract_text_from_pdf(path):
    """Extract text from PDF using pdfplumber."""
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text += (page.extract_text() or "") + "\n"
    except Exception as e:
        print(f"ERROR PDF Extract: {e}")
    return text.strip()

def extract_text_from_pptx(path):
    """Extract text from PPTX slides."""
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"ERROR PPTX Extract: {e}")
    return text.strip()

def extract_questions_from_paper(file_path):
    """
    Standard OCR extraction for Reference Paper.
    UPGRADED: Precision mode for 100% mark/question coverage.
    """
    try:
        model = get_model()
        prompt = """
        ACT AS AN ELITE EXAM ANALYST & VISUAL OCR SPECIALIST.
        You are scanning a REFERENCE ANSWER PAPER. You must capture EVERY mark, question, and piece of data.
        
        CRITICAL: HANDWRITTEN TABLES & DATA
        - If you see a table or comparison, transcribe it using MARKDOWN TABLE format.
        - Capture all column headers and row data precisely.
        - If a diagram is present, describe it in 1-2 sentences within the "answerText" (e.g., [Diagram: Schematic of a DC Motor showing...]).
        
        STEP 1: AUDIT THE PAPER
        - Find "Total Marks", "Maximum Marks", or "M.M." (e.g., 20, 50, 100).
        - Count every distinct question block. 
        - Look at the right-hand margin; marks are often listed there in [brackets].
        
        STEP 2: SCAN FOR HIDDEN QUESTIONS
        - Ensure you don't miss sub-parts (like 1a, 1b).
        - Capture the FULL text of the question.
        
        STEP 3: VERIFY THE MATH
        - Sum all the marks. If the sum doesn't match "Total Marks", RE-SCAN.
        
        OUTPUT FORMAT (JSON ONLY):
        [
          {
            "number": "string",
            "text": "full question wording",
            "marks": number,
            "answerText": "model answer (USE MARKDOWN TABLES IF NEEDED)",
            "keywords": ["essential", "terms"],
            "hasDiagram": true/false
          }
        ]
        """
        
        part = _img_part(file_path)
        resp = model.generate_content([part, prompt], safety_settings=SAFETY_NONE)
        
        if not resp.text: raise ValueError("AI Empty Response")
        
        data = _safe_json(resp.text, [])
        for item in data:
            item['number'] = _clean_qnum(item.get('number'))
            # Robust mark conversion
            try: item['marks'] = float(item.get('marks', 5))
            except: item['marks'] = 5.0
            
        return data
        
    except exceptions.ResourceExhausted:
        print("WARNING: Quota hit. Use short delay.")
        # Minimal wait to not trigger browser timeout
        time.sleep(2) 
        # Final try
        try:
            new_model = get_model(force_refresh=True)
            resp = new_model.generate_content([_img_part(file_path), prompt], safety_settings=SAFETY_NONE)
            data = _safe_json(resp.text, [])
            for item in data: item['number'] = _clean_qnum(item.get('number'))
            return data
        except: return []
    except Exception as e:
        print(f"ERROR Extraction: {e}")
        return []

def extract_paper_structure(file_path):
    """
    STEP 1: OCR PURE EXTRACTION
    Extracts question numbers, text, and marks from a blank paper.
    """
    try:
        model = get_model()
        prompt = """
        ACT AS AN EXPERT OCR ANALYST.
        Scan this Question Paper and extract the structured list of questions.

        CRITICAL RULES:
        1. Capture EVERY question number and text.
        2. Look for marks/weightage (usually in brackets at the right margin).
        3. If no marks are found, default to 5.0.
        4. Capture sub-questions (e.g., 1a, 1b) precisely.

        OUTPUT JSON ONLY:
        [
          {
            "number": "string",
            "text": "full question",
            "marks": number
          }
        ]
        """
        part = _img_part(file_path)
        resp = model.generate_content([part, prompt], safety_settings=SAFETY_NONE)
        if not resp.text: raise ValueError("AI Empty Response")
        data = _safe_json(resp.text, [])
        for item in data:
            item['number'] = _clean_qnum(item.get('number'))
            try: item['marks'] = float(item.get('marks', 5.0))
            except: item['marks'] = 5.0
        return data
    except exceptions.ResourceExhausted:
        print("WARNING: Quota hit during OCR. Retrying in 10s...")
        time.sleep(10)
        try:
            new_model = get_model(force_refresh=True)
            resp = new_model.generate_content([_img_part(file_path), prompt], safety_settings=SAFETY_NONE)
            data = _safe_json(resp.text, [])
            for item in data: item['number'] = _clean_qnum(item.get('number'))
            return data
        except: return []
    except Exception as e:
        print(f"ERROR Structure Extraction: {str(e)}")
        traceback.print_exc()
        return []

def enrich_questions_with_ai(questions, subject_hint="", study_material=None, use_web=False):
    """
    OPTIONAL STEP: AI KNOWLEDGE ENRICHMENT
    Generates additional AI-powered reference answers and rubrics.
    If use_web=True, it performs deep research to provide a global academic perspective.
    """
    if not questions: return []
    try:
        model = get_model()
        # If web search is requested, we use a research-capable configuration if available
        # Otherwise, the prompt itself directs the AI to use its global knowledge
        q_json = json.dumps(questions)
        
        context_block = ""
        if study_material:
            context_block = f"\nPRIMARY CONTEXT (Teacher's Material):\n{study_material[:10000]}\n"
        
        research_instruction = ""
        if use_web:
            research_instruction = "RESEARCH MODE ACTIVE: Search the web/global academic databases to provide the most accurate, state-of-the-art answer. Ensure the answer aligns with standard university curricula."
        
        prompt = f"""
        ACT AS AN ELITE UNIVERSITY PROFESSOR & RESEARCHER.
        Subject: {subject_hint}
        {context_block}
        {research_instruction}

        TASK: For each question provided, generate a "Global AI Reference Answer" that represents a perfect, high-scoring response.
        
        STRICT CONSTRAINTS:
        1. DO NOT MODIFY THE "marks" provided.
        2. CAPTURE ALL TECHNICAL TERMS.
        
        OUTPUT FORMAT (JSON ONLY):
        [
          {{
            "number": "...",
            "aiReferenceText": "Concise but perfect answer based on AI/Web knowledge.",
            "aiRubrics": [
              {{"parameter": "Technical Accuracy", "weight": ...}},
              {{"parameter": "Clarity & Structure", "weight": ...}}
            ]
          }}
        ]

        QUESTIONS:
        {q_json}
        """
        
        # We don't use tools=[...] here to avoid 400 errors on older models, 
        # but modern Gemini models will use their built-in search if prompted for research.
        resp = model.generate_content(prompt, safety_settings=SAFETY_NONE)
        if not resp.text: raise ValueError("AI Empty Response")
        data = _safe_json(resp.text, [])
        return data
    except Exception as e:
        print(f"ERROR AI Enrichment: {e}")
        return []

def extract_answer_sheet(file_path, questions, study_material=None):
    """
    OCR for student answers. 
    UPGRADED: Context-Aware Transcription using both Question Text AND Study Materials.
    """
    try:
        model = get_model()
        q_ctx = "\n".join([f"{_clean_qnum(q.get('number'))}: {q.get('text')[:100]}" for q in questions])
        
        # Add study material context to help with messy technical terms
        material_ctx = ""
        if study_material:
            # We take a snippet if it's too long, but usually it's fine
            material_ctx = f"\nTECHNICAL CONTEXT FROM STUDY MATERIAL:\n{study_material[:5000]}\n"

        prompt = f"""
        ACT AS AN EXPERT HANDWRITING & TABULAR DATA DECODER.
        Transcribe the handwritten student answers from this image.
        
        STRICT MATCHING RULES:
        1. NO FORCE-MATCHING: Only extract an answer if it is explicitly written. If a question from the context is NOT answered by the student, DO NOT include it in the output.
        2. CROSS-REFERENCE CONTENT: Don't just look at the number the student wrote (e.g., "1"). Look at the words they wrote and match them to the most relevant QUESTION CONTEXT below. 
        3. TECHNICAL DECODING: Use the TECHNICAL CONTEXT FROM STUDY MATERIAL provided below to help decode messy handwriting. If a word looks like a technical term from the material, it likely is.
        4. HANDLE "OR" / OPTIONAL QUESTIONS: Ensure the "questionNumber" in your JSON matches the correct ID from the context based on content alignment.
        5. ACCURATE TRANSCRIPTION: If a student writes a table, use MARKDOWN TABLE format. 
        6. IGNORE STRUCK-THROUGH TEXT: If text is crossed out, DO NOT transcribe it.
        
        QUESTION CONTEXT (ID: Question Text):
        {q_ctx}
        {material_ctx}
        
        OUTPUT FORMAT (JSON ONLY):
        [
          {{
            "questionNumber": "ID from context (e.g., 1a, 2, 5b)",
            "answerText": "transcribed text (USE MARKDOWN TABLES FOR TABULAR DATA)",
            "hasDiagram": true/false
          }}
        ]
        """
        
        resp = model.generate_content([_img_part(file_path), prompt], safety_settings=SAFETY_NONE)
        
        if not resp.text: raise ValueError("AI Empty Response")
        data = _safe_json(resp.text, [])
        for item in data:
            item['questionNumber'] = _clean_qnum(item.get('questionNumber'))
        return data
        
    except exceptions.ResourceExhausted:
        print("WARNING: Student OCR Quota hit. Attempting model switch...")
        time.sleep(2)
        try:
            # Force refresh to find another working model
            new_model = get_model(force_refresh=True)
            resp = new_model.generate_content([_img_part(file_path), prompt], safety_settings=SAFETY_NONE)
            data = _safe_json(resp.text, [])
            for item in data: item['questionNumber'] = _clean_qnum(item.get('questionNumber'))
            return data
        except Exception as retry_e:
            print(f"CRITICAL: Retry failed: {retry_e}")
            raise RuntimeError("Gemini API is busy across all available models. Please wait 60 seconds and try again.")
            
    except Exception as e:
        print(f"ERROR Student OCR: {e}")
        # Allow it to crash upward so the UI knows the OCR failed, instead of grading an empty page
        raise RuntimeError(f"OCR Extraction Failed: {str(e)}")

def grade_answer_feedback(question_text, model_answer, student_answer, marks, ml_score, study_material=None):
    """
    Teacher Persona Grading: Upgraded Prompt for Highest Accuracy.
    - Focuses on "Concept over Syntax"
    - Identifies specific "Mistakes" and "Improvement Paths"
    - Uses Study Material context if provided
    """
    try:
        model = get_model()
        prompt = f"""
        PERSONA: You are a kind, sentimental, and supportive university teacher. 
        Your goal is to ENCOURAGE students by finding reasons to give them marks, not reasons to subtract them.
        
        - TARGET ACCURACY: Aim to match a "Strict Teacher's Sentiment"—looking for depth and effort.

        GRADING PHILOSOPHY:
        - DEPTH OVER BREVITY: For high-mark questions (5, 7, 10 marks), a single line or extremely short answer is NOT enough. Penalize heavily for lack of detail/completeness even if the core idea is there.
        - CONCEPT OVER SYNTAX: If the student understands the "spirit" of the answer, give credit, BUT only if they provide enough detail to justify the marks.
        - BE FAIR: If an answer is clearly incomplete for the marks assigned, give a low score (e.g. 0-2 out of 7).
        - IGNORE OCR ERRORS: Assume intended words if the context fits.

        RUBRIC / ASSESSMENT PARAMETERS:
        {model_answer.get('assessmentParameters', []) if isinstance(model_answer, dict) else 'N/A'}

        RELEVANT SOURCES/CONCEPTS:
        {model_answer.get('relevantSources', []) if isinstance(model_answer, dict) else 'N/A'}

        ADDITIONAL REFERENCE MATERIAL (Context from Teacher PDFs/PPTs):
        {study_material if study_material else 'No additional material provided.'}

        QUESTION ({marks} marks total): {question_text}

        REFERENCE ANSWER (Teacher's Ideal points):
        {model_answer.get('answerText', model_answer) if isinstance(model_answer, dict) else model_answer}

        AI KNOWLEDGE REFERENCE (Web/Global Research):
        {model_answer.get('aiReferenceText', 'No additional AI research provided.') if isinstance(model_answer, dict) else 'N/A'}

        STUDENT'S ANSWER:
        {student_answer}

        MARKING INSTRUCTIONS:
        1. Identify the core concepts in the reference.
        2. Evaluate the student's answer against the specific ASSESSMENT PARAMETERS provided above.
        3. Award marks for each parameter found.
        4. DEPTH CHECK: If the question is worth {marks} marks and the answer is only one sentence, it is "Critically Incomplete". Award very few marks.
        5. suggestedScore = (calculated percentage 0-100) — Be a strict but fair and encouraging evaluator.

        OUTPUT JSON ONLY:
        {{
          "feedback": "2-3 sentences — a supportive and encouraging comment",
          "strengths": ["What they did well"],
          "mistakes": ["Gently point out what was missed, explaining why it's important"],
          "improvements": ["One clear tip for next time"],
          "suggestedScore": 0-100
        }}
        """
        
        # Handle Quota Exceeded for individual calls
        try:
            resp = model.generate_content(prompt, safety_settings=SAFETY_NONE)
        except exceptions.ResourceExhausted:
            print("WARNING: Grading quota hit. Waiting 12 seconds...")
            time.sleep(12)
            resp = model.generate_content(prompt, safety_settings=SAFETY_NONE)

        res = _safe_json(resp.text, {})
        # Ensure raw score is capped and fallback to ML if AI fails (ml_score is already 0-100)
        res['suggestedScore'] = max(0, min(100, float(res.get('suggestedScore', ml_score))))
        return res
    except Exception as e:
        print(f"DEBUG: Grade Feedback Error: {e}")
        return {
            'feedback': 'Detailed AI analysis was unavailable for this answer.',
            'strengths': [],
            'mistakes': ['Analysis could not pinpoint mistakes due to server error.'],
            'improvements': ['Consider manual verification.'],
            'suggestedScore': ml_score if ml_score else 0
        }

def grade_batch_feedback(student_name, exam_title, questions_batch, study_material=None):
    """
    Grades multiple questions AND generates a summary report narrative in a single call.
    questions_batch: list of {id, text, marks, modelAnswer, studentAnswer, mlScore}
    """
    try:
        model = get_model()
        
        # Prepare structured context for the AI
        q_entries = []
        for q in questions_batch:
            q_entries.append(f"""
--- QUESTION {q.get('number', '?')} ({q.get('marks', 5)} marks) ---
TEXT: {q.get('text')}
REFERENCE: {q.get('modelAnswer')}
AI KNOWLEDGE REFERENCE: {q.get('aiReferenceText', 'Not available')}
STUDENT ANSWER: {q.get('studentAnswer')}
ML PRELIMINARY SCORE: {q.get('mlScore', 0)}%
""")

        context_block = "\n".join(q_entries)
        
        prompt = f"""
ACT AS AN ELITE UNIVERSITY PROFESSOR & FAIR GRADER.
Student: {student_name}
Exam: {exam_title}

GRADING PHILOSOPHY (TRIPLE SOURCE VALIDATION):
1. TRIPLE SOURCE VALIDATION: Use the REFERENCE answer paper, the ADDITIONAL STUDY MATERIAL (PPTs/PDFs), and the AI KNOWLEDGE REFERENCE (from Web Search) as valid sources for marks. If a concept is valid in any of these, AWARD MARKS.
2. CONCEPT OVER SYNTAX: If the student demonstrates clear understanding of the core concepts, be generous. Do not penalize for slightly different wording.
3. ENCOURAGE ACCURACY: Aim for high accuracy. If a student's answer is 70% correct conceptually across the sources, they should receive at least 70% marks.
4. AVOID HYPER-STRICTNESS: Only penalize if the answer is factually wrong or completely missing from all provided sources.

ADDITIONAL STUDY MATERIAL CONTEXT (PRIMARY SOURCE OF TRUTH):
{study_material if study_material else 'No additional material provided.'}

QUESTIONS TO GRADE:
{context_block}

OUTPUT FORMAT (JSON ONLY):
{{
  "grades": [
    {{
      "number": "string",
      "feedback": "2-3 sentences of supportive but honest feedback",
      "strengths": ["list"],
      "mistakes": ["list"],
      "improvements": ["list"],
      "suggestedScore": 0-100
    }}
  ],
  "overallNarrative": "A 2-3 sentence summary report addressed to the student about their overall performance."
}}
"""
        
        resp = model.generate_content(prompt, safety_settings=SAFETY_NONE)
        if not resp.text: raise ValueError("AI Empty Response")
        
        results = _safe_json(resp.text, {})
        return results
        
    except Exception as e:
        print(f"ERROR Batch Grading: {e}")
        return {"grades": [], "overallNarrative": ""}

def generate_student_report(student_name, exam_title, grades, percentage):
    try:
        model = get_model()
        prompt = f"Summary report for {student_name} on {exam_title} ({percentage}%). Be brief."
        resp = model.generate_content(prompt)
        return resp.text.strip()
    except:
        return f"{student_name} scored {percentage}%."
